from cgitb import text
from email.mime import image
from fileinput import filename
from re import search
import fitz
from PIL import Image
from pptx import Presentation
from pptx.util import Pt
import os
import io
from io import BytesIO


posDict = {
    # Define coordinates for positioning. Format is test name then test type. e.g. a VT-07 test with a mediumwave plot
    # Tuple Order: Left , Top, Width, Height
    # VT-07
    'VT07MW' : (Pt(1), Pt(70), Pt(233), Pt(176)),
    'VT07FM1' : (Pt(239), Pt(70), Pt(233), Pt(176)),
    'VT07FM2' : (Pt(477), Pt(70), Pt(233), Pt(176)),
    'VT07DAB1AV' : (Pt(1), Pt(275), Pt(175), Pt(130)),
    'VT07DAB1RMS' : (Pt(180), Pt(275), Pt(175), Pt(130)),
    'VT07DAB2AV' : (Pt(360), Pt(275), Pt(175), Pt(130)),
    'VT07DAB2RMS' : (Pt(540), Pt(275), Pt(175), Pt(130)),
    # VT-01 3 Metre
    'VT01ThreeVertical' : (Pt(1), Pt(85), Pt(358), Pt(270)),
    'VT01ThreeHorizontal' : (Pt(360), Pt(85), Pt(358), Pt(270)),
    # VT-12 Single Phase
    'VT12SingleL1' : (Pt(1), Pt(85), Pt(358), Pt(270)),
    'VT12SingleN' : (Pt(360), Pt(85), Pt(358), Pt(270)),
    # VT-12 Three Phase
    'VT12TripleL1' : (Pt(1), Pt(65), Pt(221), Pt(167)),
    'VT12TripleL2' : (Pt(222), Pt(65), Pt(221), Pt(167)),
    'VT12TripleL3' : (Pt(1), Pt(240), Pt(221), Pt(167)),
    'VT12TripleN' : (Pt(222), Pt(240), Pt(221), Pt(167)),
    # VT-15 Electric
    'VT15E16' : (Pt(1), Pt(85), Pt(233), Pt(176)),
    'VT15E40' : (Pt(239), Pt(85), Pt(233), Pt(176)),
    'VT15E70' : (Pt(477), Pt(85), Pt(233), Pt(176)),
    # VT-15 Magnetic Radial and Transverse
    'VT15HR16' : (Pt(1), Pt(65), Pt(221), Pt(167)),
    'VT15HR40' : (Pt(239), Pt(65), Pt(221), Pt(167)),
    'VT15HR70' : (Pt(477), Pt(65), Pt(221), Pt(167)),
    'VT15HT16' : (Pt(1), Pt(235), Pt(221), Pt(167)),
    'VT15HT40' : (Pt(239), Pt(235), Pt(221), Pt(167)),
    'VT15HT70' : (Pt(477), Pt(235), Pt(221), Pt(167)),
}

cropDict = {
    # Define coordinates for cropping. Old refers to the old style PDFs we use that crop in two static positions
    # Left Start, Top Start, Left End, Top End
    'upperOld' : ((130, 138, 1000, 800)),
    'lowerOld' : ((130, 820, 1000, 1482)),
    'upperOldMagnetic' : ((130, 270, 1000, 932))
}

nameDict = {
    # Takes names of functions and writes their proper name onto the slide
    'VT01Ten' : 'VT-01 Off-Board Emissions (10m)',
    'VT01Three' : 'VT-01 Off-Board Emissions (3m)',
    'VT07' : 'VT-07 On-Board Emissions',
    'VT12Single' : 'VT-12 Conducted Emissions (Single Phase)',
    'VT12Triple' : 'VT-12 Conducted Emissions (Three Phase)',
    'VT15Electric' : 'VT-15 Electric Fields',
    'VT15Magnetic' : 'VT-15 Magnetic Fields'
}

# Defining global variables
extractedImages = []
croppedImages = []
rejectedList = []
slideCounter = 0


def searchReplace(search_str, repl_str, input, output):
    # Attempts to search and replace on the entire file. Likely needs rewriting to be more robust and not need a template
    prs = Presentation(input)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                if(shape.text.find(search_str))!=-1:
                    text_frame = shape.text_frame
                    cur_text = text_frame.paragraphs[0].runs[0].text
                    new_text = cur_text.replace(str(search_str), str(repl_str))
                    text_frame.paragraphs[0].runs[0].text = new_text
    prs.save(output)
    print(search_str + ' replaced with ' + repl_str)


def extractImages(PDFName, image_folder):
    fileName = image_folder + '/' + PDFName
    doc = fitz.open(fileName)
    zoom = 2 # to increase the resolution
    mat = fitz.Matrix(zoom, zoom)
    noOfPages = doc.pageCount
    for pageNo in range(noOfPages):
        page = doc.load_page(pageNo) #number of page
        pix = page.get_pixmap(matrix = mat)
        extractedImages.append(pix)
        print('Converting PDFs to Image')


def cropGraph(targetImg, cropTuple, imName):
    targetPIL = targetImg.tobytes("PNG")
    im = Image.open(io.BytesIO(targetPIL))
    im1 = im.crop(box=cropTuple)
    croppedImages.append(im1)
    print(imName +  ' cropped')


def insertImage(oldFileName, newFileName, img, positionTuple, slideNumber):
    # Inserts an image from the croppedImages array into slideNumber using a position from posDict
    prs = Presentation(oldFileName)
    slide = prs.slides[slideNumber]
    left = positionTuple[0]
    top = positionTuple[1]
    width = positionTuple[2]
    height = positionTuple[3]
    temp = BytesIO()
    img.save(temp, "PNG")
    slide.shapes.add_picture(temp, left, top, width, height)
    prs.save(newFileName)
    print('Image inserted')


def initialisePowerPoint(emptyDeckName, newDeckName):
    # Sets up the empty, fresh PPTX file
    emptyDeckName = emptyDeckName + '.pptx'
    newDeckName = newDeckName + '.pptx'
    prs = Presentation(emptyDeckName)
    prs.save(newDeckName)
    print('Created new PowerPoint: ' + newDeckName)



def VT07(PDFName, folderName, slideNumber, deckName):
    extractImages(PDFName, folderName)
    cropGraph(extractedImages[1], cropDict['upperOld'], 'MW')
    cropGraph(extractedImages[1], cropDict['lowerOld'], 'FM1')
    cropGraph(extractedImages[2], cropDict['upperOld'], 'FM2')
    cropGraph(extractedImages[2], cropDict['lowerOld'], 'DAB1AV')
    cropGraph(extractedImages[3], cropDict['upperOld'], 'DAB1RMS')
    cropGraph(extractedImages[3], cropDict['lowerOld'], 'DAB2AV')
    cropGraph(extractedImages[4], cropDict['upperOld'], 'DAB2RMS')
    deckName = deckName + '.pptx'
    insertImage(deckName, deckName, croppedImages[0], posDict['VT07MW'], slideNumber)
    insertImage(deckName, deckName, croppedImages[1], posDict['VT07FM1'], slideNumber) 
    insertImage(deckName, deckName, croppedImages[2], posDict['VT07FM2'], slideNumber) 
    insertImage(deckName, deckName, croppedImages[3], posDict['VT07DAB1AV'], slideNumber) 
    insertImage(deckName, deckName, croppedImages[4], posDict['VT07DAB1RMS'], slideNumber) 
    insertImage(deckName, deckName, croppedImages[5], posDict['VT07DAB2AV'], slideNumber) 
    insertImage(deckName, deckName, croppedImages[6], posDict['VT07DAB2RMS'], slideNumber)
    extractedImages.clear()
    croppedImages.clear()
    print('Finished VT-07 for ' + PDFName)

def VT01Three(PDFName, folderName, slideNumber, deckName):
    extractImages(PDFName, folderName)
    cropGraph(extractedImages[1], cropDict['upperOld'], 'VT01ThreeVertical')
    cropGraph(extractedImages[1], cropDict['lowerOld'], 'VT01ThreeHorizontal')
    deckName = deckName + '.pptx'
    insertImage(deckName, deckName, croppedImages[0], posDict['VT01ThreeVertical'], slideNumber)
    insertImage(deckName, deckName, croppedImages[1], posDict['VT01ThreeHorizontal'], slideNumber)
    extractedImages.clear()
    croppedImages.clear()
    print('Finished VT-01 3m for ' + PDFName)

def VT12Single(PDFName, folderName, slideNumber, deckName):
    extractImages(PDFName, folderName)
    cropGraph(extractedImages[1], cropDict['upperOld'], 'VT12SingleL1')
    cropGraph(extractedImages[1], cropDict['lowerOld'], 'VT12SingleN')
    deckName = deckName + '.pptx'
    insertImage(deckName, deckName, croppedImages[0], posDict['VT12SingleL1'], slideNumber)
    insertImage(deckName, deckName, croppedImages[1], posDict['VT12SingleN'], slideNumber)
    extractedImages.clear()
    croppedImages.clear()
    print('Finished VT-12 Single Phase for ' + PDFName)

def VT12Triple(PDFName, folderName, slideNumber, deckName):
    extractImages(PDFName, folderName)
    cropGraph(extractedImages[1], cropDict['upperOld'], 'VT12TripleL1')
    cropGraph(extractedImages[1], cropDict['lowerOld'], 'VT12TripleL2')
    cropGraph(extractedImages[2], cropDict['upperOld'], 'VT12TripleL3')
    cropGraph(extractedImages[2], cropDict['lowerOld'], 'VT12TripleN')
    deckName = deckName + '.pptx'
    insertImage(deckName, deckName, croppedImages[0], posDict['VT12TripleL1'], slideNumber)
    insertImage(deckName, deckName, croppedImages[1], posDict['VT12TripleL2'], slideNumber)
    insertImage(deckName, deckName, croppedImages[2], posDict['VT12TripleL3'], slideNumber)
    insertImage(deckName, deckName, croppedImages[3], posDict['VT12TripleN'], slideNumber)
    extractedImages.clear()
    croppedImages.clear()
    print('Finished VT-12 Three Phase for ' + PDFName)

def VT15Electric(PDFName, folderName, slideNumber, deckName):
    extractImages(PDFName, folderName)
    cropGraph(extractedImages[1], cropDict['upperOld'], 'VT15E16')
    cropGraph(extractedImages[1], cropDict['lowerOld'], 'VT15E40')
    cropGraph(extractedImages[2], cropDict['upperOld'], 'VT15E70')
    deckName = deckName + '.pptx'
    insertImage(deckName, deckName, croppedImages[0], posDict['VT15E16'], slideNumber)
    insertImage(deckName, deckName, croppedImages[1], posDict['VT15E40'], slideNumber)
    insertImage(deckName, deckName, croppedImages[2], posDict['VT15E70'], slideNumber)
    extractedImages.clear()
    croppedImages.clear()
    print('Finished VT-15 Electric Field for ' + PDFName)

def VT15Magnetic(PDFName, folderName, slideNumber, deckName):
    extractImages(PDFName, folderName)
    cropGraph(extractedImages[1], cropDict['upperOldMagnetic'], 'VT15HR16')
    cropGraph(extractedImages[2], cropDict['upperOld'], 'VT15HR40')
    cropGraph(extractedImages[2], cropDict['lowerOld'], 'VT15HR70')
    cropGraph(extractedImages[3], cropDict['upperOld'], 'VT15HT16')
    cropGraph(extractedImages[3], cropDict['lowerOld'], 'VT15HT40')
    cropGraph(extractedImages[4], cropDict['upperOld'], 'VT15HT70')
    deckName = deckName + '.pptx'
    insertImage(deckName, deckName, croppedImages[0], posDict['VT15HR16'], slideNumber)
    insertImage(deckName, deckName, croppedImages[1], posDict['VT15HR40'], slideNumber)
    insertImage(deckName, deckName, croppedImages[2], posDict['VT15HR70'], slideNumber)
    insertImage(deckName, deckName, croppedImages[3], posDict['VT15HT16'], slideNumber)
    insertImage(deckName, deckName, croppedImages[4], posDict['VT15HT40'], slideNumber)
    insertImage(deckName, deckName, croppedImages[5], posDict['VT15HT70'], slideNumber)
    extractedImages.clear()
    croppedImages.clear()
    print('Finished VT-15 Electric Field for ' + PDFName)

def setSlideCounter(num):
    global slideCounter
    slideCounter = num
    print('Slide counter set to ' + str(slideCounter))


def loopFolder(folderName, deckName, reportFunction):

    directory = folderName
    global slideCounter
    
    for file in os.listdir(directory):
        if file.endswith(".Pdf") or file.endswith(".pdf"):
            print('Working on slide ' + str(slideCounter) + ', File Name: ' + file)
            reportFunction(file, folderName, slideCounter, deckName)
            searchString = '*' + str(slideCounter) + '*'
            replaceString = str(file)[:-4] + ' | ' + (nameDict[str(reportFunction.__name__)])
            searchReplace(searchString, replaceString, deckName + '.pptx', deckName + '.pptx')
            slideCounter = slideCounter + 1
    print('Finished with folder: ' + folderName)





initialisePowerPoint('emptyDeck', 'newDeck')

setSlideCounter(0)

#loopFolder('VT-01 3m','newDeck', VT01Three)
#loopFolder('VT-07','newDeck', VT07)
#loopFolder('VT-12 Single Phase', 'newDeck', VT12Single)
#loopFolder('VT-12 Three Phase', 'newDeck', VT12Triple)
#loopFolder('VT-15 Electric', 'newDeck', VT15Electric)
loopFolder('VT-15 Magnetic', 'newDeck', VT15Magnetic)

print('Finished all jobs...')
